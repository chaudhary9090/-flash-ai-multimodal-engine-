"""
Pytest Suite: Integration & PPTX Document Ingestion Tests
"""

import io
import pytest
from fastapi.testclient import TestClient
from app.main import app
from app.services.document_service import extract_clean_text

client = TestClient(app)


def test_health_endpoint():
    response = client.get("/api/v1/health")
    assert response.status_code == 200
    data = response.json()
    assert data["status"] == "healthy"


def test_auth_registration_and_login():
    reg_res = client.post("/api/v1/auth/register", json={"username": "demo_user", "password": "password123"})
    assert reg_res.status_code in [200, 201]

    login_res = client.post("/api/v1/auth/login", json={"username": "demo_user", "password": "password123"})
    assert login_res.status_code == 200
    assert "access_token" in login_res.json()


def test_pptx_clean_text_extraction():
    """Verifies that PPTX files do NOT output binary ZIP garbage (PK signatures)."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "The Kindness Quest Presentation"
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Empathy and Positive Social Impact"

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_bytes = pptx_io.getvalue()

    # Extract clean text
    extracted_text = extract_clean_text("test_presentation.pptx", pptx_bytes)
    
    # Assert clean text without PK ZIP headers
    assert not extracted_text.startswith("PK")
    assert "The Kindness Quest Presentation" in extracted_text
    assert "Empathy and Positive Social Impact" in extracted_text


def test_unsupported_file_type_handling():
    files = {"file": ("malicious.exe", b"MZ\x90\x00\x03\x00\x00\x00", "application/octet-stream")}
    response = client.post("/api/v1/upload", files=files)
    assert response.status_code == 415
    data = response.json()
    assert "Unsupported file type" in data["detail"]
